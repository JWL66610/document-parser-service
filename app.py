import io
import os
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import httpx
from docx import Document
from fastapi import FastAPI, Header, HTTPException, Request
from pydantic import BaseModel, Field, HttpUrl, ValidationError
from pypdf import PdfReader
from pptx import Presentation

MAX_DOWNLOAD_BYTES = 20 * 1024 * 1024
DEFAULT_MAX_CHARS = 12_000
MAX_TEXT_CHARS = 30_000
MAX_REDIRECTS = 3

app = FastAPI(title="Xuetuzhiban Document Parser", version="1.0.0")


class ExtractRequest(BaseModel):
    file_url: HttpUrl
    file_name: str = Field(min_length=1, max_length=255)
    max_chars: int = Field(default=DEFAULT_MAX_CHARS, ge=500, le=MAX_TEXT_CHARS)


class ExtractResponse(BaseModel):
    filename: str
    document_type: str
    text: str
    truncated: bool


def allowed_hosts() -> set[str]:
    configured = os.getenv("ALLOWED_FILE_HOSTS", "ai.yznu.edu.cn")
    return {host.strip().lower() for host in configured.split(",") if host.strip()}


def require_api_key(api_key: str | None) -> None:
    expected = os.getenv("PARSER_API_KEY")
    if not expected:
        raise HTTPException(status_code=503, detail="Parser service is not configured")
    if api_key != expected:
        raise HTTPException(status_code=401, detail="Invalid API key")


def validate_url(url: str) -> None:
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"}:
        raise HTTPException(status_code=400, detail="Only HTTP(S) file URLs are allowed")
    if parsed.hostname not in allowed_hosts():
        raise HTTPException(status_code=400, detail="File URL host is not allowed")


async def download_file(url: str) -> bytes:
    try:
        async with httpx.AsyncClient(
            follow_redirects=False,
            timeout=httpx.Timeout(30.0, connect=10.0),
        ) as client:
            next_url = url
            for _ in range(MAX_REDIRECTS + 1):
                validate_url(next_url)

                async with client.stream("GET", next_url) as response:
                    if response.is_redirect:
                        location = response.headers.get("location")
                        if not location:
                            raise HTTPException(
                                status_code=502,
                                detail="File download redirect has no location",
                            )
                        next_url = str(response.url.join(location))
                        validate_url(next_url)
                        continue

                    response.raise_for_status()
                    content_length = response.headers.get("content-length")
                    if content_length and int(content_length) > MAX_DOWNLOAD_BYTES:
                        raise HTTPException(
                            status_code=413,
                            detail="File exceeds the 20 MB parsing limit",
                        )

                    chunks: list[bytes] = []
                    size = 0
                    async for chunk in response.aiter_bytes():
                        size += len(chunk)
                        if size > MAX_DOWNLOAD_BYTES:
                            raise HTTPException(
                                status_code=413,
                                detail="File exceeds the 20 MB parsing limit",
                            )
                        chunks.append(chunk)
                    return b"".join(chunks)

            raise HTTPException(status_code=502, detail="Too many file download redirects")
    except HTTPException:
        raise
    except httpx.HTTPError as error:
        raise HTTPException(status_code=502, detail=f"Unable to download file: {error}") from error


def extract_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_docx(file_bytes: bytes) -> str:
    document = Document(io.BytesIO(file_bytes))
    parts = [paragraph.text for paragraph in document.paragraphs]

    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))

    return "\n".join(parts)


def extract_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(file_bytes))
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"Slide {index}:\n" + "\n".join(texts))

    return "\n\n".join(slides)


def extract_txt(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "gb18030", "gbk"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise HTTPException(status_code=422, detail="TXT encoding could not be detected")


def extract_text(file_name: str, file_bytes: bytes) -> tuple[str, str]:
    extension = Path(file_name).suffix.lower()

    try:
        if extension == ".pdf":
            return extract_pdf(file_bytes), "pdf"
        if extension == ".docx":
            return extract_docx(file_bytes), "docx"
        if extension == ".pptx":
            return extract_pptx(file_bytes), "pptx"
        if extension == ".txt":
            return extract_txt(file_bytes), "txt"
    except Exception as error:
        raise HTTPException(status_code=422, detail=f"Document parsing failed: {error}") from error

    if extension in {".doc", ".ppt"}:
        raise HTTPException(
            status_code=415,
            detail="Legacy DOC/PPT is not supported. Convert it to DOCX/PPTX first.",
        )
    raise HTTPException(status_code=415, detail="Only PDF, DOCX, PPTX, and TXT are supported")


def normalize_text(text: str) -> str:
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


async def read_extract_request(request: Request) -> ExtractRequest:
    content_type = request.headers.get("content-type", "").lower()

    try:
        if content_type.startswith("application/json"):
            payload: dict[str, Any] = await request.json()
        elif content_type.startswith(
            ("multipart/form-data", "application/x-www-form-urlencoded")
        ):
            form_data = await request.form()
            payload = {
                "file_url": form_data.get("file_url"),
                "file_name": form_data.get("file_name"),
                "max_chars": form_data.get("max_chars", DEFAULT_MAX_CHARS),
            }
        else:
            raise HTTPException(
                status_code=415,
                detail="Use JSON, form-data, or x-www-form-urlencoded request body",
            )

        return ExtractRequest.model_validate(payload)
    except HTTPException:
        raise
    except (ValidationError, ValueError, TypeError) as error:
        raise HTTPException(status_code=422, detail=f"Invalid request body: {error}") from error


@app.get("/health")
async def health() -> dict[str, str]:
    return {"status": "ok"}


@app.post("/v1/extract", response_model=ExtractResponse)
async def extract_document(
    request: Request,
    x_api_key: str | None = Header(default=None),
) -> ExtractResponse:
    require_api_key(x_api_key)

    extract_request = await read_extract_request(request)
    file_bytes = await download_file(str(extract_request.file_url))
    text, document_type = extract_text(extract_request.file_name, file_bytes)
    normalized = normalize_text(text)

    if not normalized:
        raise HTTPException(
            status_code=422,
            detail="No extractable text was found. Scanned PDFs require OCR.",
        )

    truncated = len(normalized) > extract_request.max_chars
    return ExtractResponse(
        filename=extract_request.file_name,
        document_type=document_type,
        text=normalized[: extract_request.max_chars],
        truncated=truncated,
    )
